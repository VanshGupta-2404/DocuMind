import os
from typing import List
from PyPDF2 import PdfReader
from docx import Document
import openpyxl
from pptx import Presentation

from embeddings import create_embeddings, add_to_vectorstore

def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from PDF file"""
    try:
        reader = PdfReader(file_path)
        text = ""
        for page_num, page in enumerate(reader.pages):
            page_text = page.extract_text()
            if page_text:
                text += f"\n--- Page {page_num + 1} ---\n{page_text}"
        return text
    except Exception as e:
        raise Exception(f"Error reading PDF: {str(e)}")

def extract_text_from_docx(file_path: str) -> str:
    """Extract text from DOCX file"""
    try:
        doc = Document(file_path)
        text = ""
        for para_num, para in enumerate(doc.paragraphs):
            if para.text.strip():
                text += f"{para.text}\n"
        return text
    except Exception as e:
        raise Exception(f"Error reading DOCX: {str(e)}")

def extract_text_from_xlsx(file_path: str) -> str:
    """Extract text from Excel file"""
    try:
        wb = openpyxl.load_workbook(file_path, data_only=True)
        text = ""
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            text += f"\n--- Sheet: {sheet_name} ---\n"
            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                if row_text.strip():
                    text += row_text + "\n"
        return text
    except Exception as e:
        raise Exception(f"Error reading Excel: {str(e)}")

def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from PowerPoint file"""
    try:
        prs = Presentation(file_path)
        text = ""
        for slide_num, slide in enumerate(prs.slides):
            text += f"\n--- Slide {slide_num + 1} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        raise Exception(f"Error reading PowerPoint: {str(e)}")

def extract_text_from_txt(file_path: str) -> str:
    """Extract text from TXT file"""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    except Exception as e:
        raise Exception(f"Error reading TXT: {str(e)}")

def extract_text(file_path: str) -> str:
    """Extract text based on file extension"""
    ext = os.path.splitext(file_path)[1].lower()
    
    extractors = {
        '.pdf': extract_text_from_pdf,
        '.docx': extract_text_from_docx,
        '.doc': extract_text_from_docx,
        '.xlsx': extract_text_from_xlsx,
        '.xls': extract_text_from_xlsx,
        '.pptx': extract_text_from_pptx,
        '.ppt': extract_text_from_pptx,
        '.txt': extract_text_from_txt,
    }
    
    extractor = extractors.get(ext)
    if not extractor:
        raise ValueError(f"Unsupported file type: {ext}")
    
    return extractor(file_path)

def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
    """Split text into overlapping chunks"""
    chunks = []
    start = 0
    text_length = len(text)
    
    while start < text_length:
        end = start + chunk_size
        chunk = text[start:end]
        
        # Try to break at sentence boundaries
        if end < text_length:
            last_period = chunk.rfind('.')
            last_newline = chunk.rfind('\n')
            break_point = max(last_period, last_newline)
            
            if break_point > chunk_size * 0.5:  # At least 50% of chunk size
                chunk = chunk[:break_point + 1]
                end = start + break_point + 1
        
        chunks.append(chunk.strip())
        start = end - overlap
    
    return [c for c in chunks if len(c) > 50]  # Filter very short chunks

def process_upload(file_path: str, filename: str) -> int:
    """Process uploaded document: extract text, chunk, and store embeddings"""
    try:
        # Extract text from document
        text = extract_text(file_path)
        
        if not text or len(text.strip()) < 10:
            raise ValueError("Document appears to be empty or unreadable")
        
        # Chunk the text
        chunks = chunk_text(text)
        
        if not chunks:
            raise ValueError("No valid text chunks created from document")
        
        # Create metadata for each chunk
        metadatas = [
            {
                "source": filename,
                "chunk_index": i,
                "total_chunks": len(chunks)
            }
            for i in range(len(chunks))
        ]
        
        # Add to vector store
        add_to_vectorstore(chunks, metadatas)
        
        return len(chunks)
    
    except Exception as e:
        raise Exception(f"Failed to process document: {str(e)}")